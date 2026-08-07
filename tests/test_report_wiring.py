"""Tests for the report-phase consumption of analysis-phase outputs.

These tests verify that after a real analysis run, the report phase:
  1. Embeds analysis figures into report.md with correct relative paths
  2. Shows per-analysis options (from the options.json files) in the document
  3. Generates image slides in the .pptx for each analysis figure
  4. Fans out multi-method analyses (cluster, dimred) into separate slides
"""

from __future__ import annotations

import json
import zipfile
from pathlib import Path

import mdtraj as md
import numpy as np
import pytest

from fastmdxplora import FastMDXplora
from fastmdxplora.cli.main import main as cli_main


# ---------------------------------------------------------------------------
# Helper: build a real trajectory + run a real analysis phase
# ---------------------------------------------------------------------------
def _make_traj_files(tmp_path: Path, *, n_residues: int = 5, n_frames: int = 25) -> Path:
    """Save a compact-globule trajectory under simulation/ so analysis can find it.

    Returns the project root (the directory containing simulation/).
    """
    rng = np.random.RandomState(0)
    top = md.Topology()
    chain = top.add_chain()
    for i in range(n_residues):
        res = top.add_residue("ALA", chain, resSeq=i + 1)
        for nm in ("N", "CA", "C", "O", "CB"):
            el = {
                "N": md.element.nitrogen,
                "CA": md.element.carbon,
                "C": md.element.carbon,
                "O": md.element.oxygen,
                "CB": md.element.carbon,
            }[nm]
            top.add_atom(nm, el, res)
    top.create_standard_bonds()

    base = rng.uniform(-0.2, 0.2, size=(top.n_atoms, 3))
    xyz = np.tile(base[None, :, :], (n_frames, 1, 1)) + rng.normal(
        scale=0.02, size=(n_frames, top.n_atoms, 3)
    )
    traj = md.Trajectory(
        xyz=xyz.astype(np.float32), topology=top, time=np.arange(n_frames) * 20.0
    )

    sim = tmp_path / "simulation"
    sim.mkdir(parents=True)
    traj[0].save_pdb(str(sim / "topology.pdb"))
    traj.save_dcd(str(sim / "production.dcd"))
    return tmp_path


@pytest.fixture
def project_with_analysis(tmp_path: Path) -> Path:
    """Project directory after a real analysis phase has run on RMSD + Rg."""
    root = _make_traj_files(tmp_path / "proj")
    fmdx = FastMDXplora(
        system=str(root / "simulation" / "topology.pdb"),
        output_dir=root,
    )
    fmdx.explore(
        include=["analysis"],
        options={"analysis": {"include": ["rmsd", "rg"]}},
    )
    return root


@pytest.fixture
def project_with_multi_method(tmp_path: Path) -> Path:
    """Project after analysis with multi-method cluster + dimred."""
    root = _make_traj_files(tmp_path / "proj", n_residues=6, n_frames=30)
    fmdx = FastMDXplora(
        system=str(root / "simulation" / "topology.pdb"),
        output_dir=root,
    )
    fmdx.explore(
        include=["analysis"],
        options={
            "analysis": {
                "include": ["cluster", "dimred"],
                "options": {
                    "cluster": {
                        "methods": ["kmeans", "hierarchical", "dbscan"],
                        "n_clusters": 3,
                    },
                    "dimred": {"methods": ["pca", "tsne"]},
                },
            }
        },
    )
    return root


# ===========================================================================
# Markdown report
# ===========================================================================
class TestReportDocument:
    def test_report_embeds_analysis_figures(self, project_with_analysis: Path):
        """The Markdown report should reference each analysis's PNG file."""
        fmdx = FastMDXplora(
            system=str(project_with_analysis / "simulation" / "topology.pdb"),
            output_dir=project_with_analysis,
        )
        fmdx.report()
        text = (project_with_analysis / "report" / "report.md").read_text(encoding="utf-8")
        # Each analysis's figure should be referenced
        assert "analysis_summary.png" in text
        assert "analysis/rmsd/rmsd.png" in text
        assert "analysis/rg/rg.png" in text
        # No "deferred" / placeholder language
        assert "deferred" not in text.lower()
        assert "wired in" not in text.lower()

    def test_report_shows_actual_options(self, project_with_analysis: Path):
        """Options from options.json should appear in the report parameters block."""
        fmdx = FastMDXplora(
            system=str(project_with_analysis / "simulation" / "topology.pdb"),
            output_dir=project_with_analysis,
        )
        fmdx.report()
        text = (project_with_analysis / "report" / "report.md").read_text(encoding="utf-8")
        # RMSD options should appear
        assert "`align`" in text
        assert "`ref`" in text
        # And the actual values, not the placeholder
        assert "Run with default options" not in text

    def test_report_includes_trajectory_metadata(self, project_with_analysis: Path):
        """The report should mention how many frames/residues were analyzed."""
        fmdx = FastMDXplora(
            system=str(project_with_analysis / "simulation" / "topology.pdb"),
            output_dir=project_with_analysis,
        )
        fmdx.report()
        text = (project_with_analysis / "report" / "report.md").read_text(encoding="utf-8")
        # We built a 25-frame, 5-residue trajectory
        assert "25 frames" in text
        assert "5 residues" in text

    def test_report_includes_all_figures_for_multi_method(
        self, project_with_multi_method: Path
    ):
        """For cluster/dimred, every method's PNG should be referenced."""
        fmdx = FastMDXplora(
            system=str(project_with_multi_method / "simulation" / "topology.pdb"),
            output_dir=project_with_multi_method,
        )
        fmdx.report()
        text = (project_with_multi_method / "report" / "report.md").read_text(encoding="utf-8")
        # cluster has two methods
        assert "cluster/cluster_kmeans.png" in text
        assert "cluster/cluster_kmeans_counts.png" in text
        assert "cluster/cluster_hierarchical.png" in text
        assert "cluster/cluster_hierarchical_counts.png" in text
        assert "cluster/cluster_hierarchical_dendrogram.png" in text
        # dimred has two methods
        assert "dimred/dimred_pca.png" in text
        assert "dimred/dimred_tsne.png" in text

    def test_report_writes_combined_analysis_summary(
        self, project_with_multi_method: Path
    ):
        fmdx = FastMDXplora(
            system=str(project_with_multi_method / "simulation" / "topology.pdb"),
            output_dir=project_with_multi_method,
        )
        fmdx.report()

        summary = project_with_multi_method / "report" / "analysis_summary.png"
        summary_manifest = (
            project_with_multi_method / "report" / "analysis_summary_manifest.json"
        )
        assert summary.is_file()
        assert summary_manifest.is_file()
        manifest = json.loads(summary_manifest.read_text(encoding="utf-8"))
        included_sources = {item["source"] for item in manifest["included"]}
        assert "analysis/cluster/cluster_kmeans_counts.png" in included_sources
        assert "analysis/cluster/cluster_hierarchical_dendrogram.png" in included_sources
        assert "analysis/dimred/dimred_pca.png" in included_sources


# ===========================================================================
# PPTX slide deck
# ===========================================================================
class TestSlideDeck:
    def test_pptx_embeds_analysis_figures(self, project_with_analysis: Path):
        """Each analysis figure should appear as an image on its own slide."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        fmdx = FastMDXplora(
            system=str(project_with_analysis / "simulation" / "topology.pdb"),
            output_dir=project_with_analysis,
        )
        fmdx.report()
        prs = Presentation(str(project_with_analysis / "report" / "slides.pptx"))

        # Count image slides
        image_slides = [
            s for s in prs.slides
            if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in s.shapes)
        ]
        # We ran RMSD + Rg, plus the combined summary slide.
        assert len(image_slides) == 3

    def test_pptx_image_slides_titled_by_analysis(self, project_with_analysis: Path):
        """Image-slide titles should be readable, not giant all-caps IDs."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        fmdx = FastMDXplora(
            system=str(project_with_analysis / "simulation" / "topology.pdb"),
            output_dir=project_with_analysis,
        )
        fmdx.report()
        prs = Presentation(str(project_with_analysis / "report" / "slides.pptx"))

        titles = {
            slide.shapes.title.text
            for slide in prs.slides
            if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
        }
        assert "Analysis summary" in titles
        assert "RMSD over frames" in titles
        assert "Radius of gyration" in titles
        assert "RG" not in titles

    def test_pptx_fans_out_multi_method_analyses(
        self, project_with_multi_method: Path
    ):
        """Multi-method analyses should produce individual restored plot slides."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        fmdx = FastMDXplora(
            system=str(project_with_multi_method / "simulation" / "topology.pdb"),
            output_dir=project_with_multi_method,
        )
        fmdx.report()
        prs = Presentation(str(project_with_multi_method / "report" / "slides.pptx"))

        image_slides = [
            s for s in prs.slides
            if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in s.shapes)
        ]
        # cluster: timelines + count plots + dendrogram, dimred: 2 methods,
        # plus the combined summary slide.
        assert len(image_slides) >= 8

    def test_pptx_multi_method_subtitle_records_method(
        self, project_with_multi_method: Path
    ):
        """Multi-method image slides should have a subtitle indicating the method."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        fmdx = FastMDXplora(
            system=str(project_with_multi_method / "simulation" / "topology.pdb"),
            output_dir=project_with_multi_method,
        )
        fmdx.report()
        prs = Presentation(str(project_with_multi_method / "report" / "slides.pptx"))

        # Find all text on image-bearing slides
        all_text = []
        for slide in prs.slides:
            if not any(
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes
            ):
                continue
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text_frame.text)
        joined = "\n".join(all_text).lower()
        # The subtitle text should include method names somewhere
        for method in ("kmeans", "hierarchical", "pca", "tsne"):
            assert method in joined, f"method {method} missing from slide text"
        assert "analysis summary" in joined


# ===========================================================================
# Static HTML dashboard
# ===========================================================================
class TestDashboard:
    def test_dashboard_is_generated_and_links_outputs(self, project_with_analysis: Path):
        rmsd_png = project_with_analysis / "analysis" / "rmsd" / "rmsd.png"
        rg_png = project_with_analysis / "analysis" / "rg" / "rg.png"
        original_rmsd = rmsd_png.read_bytes()
        original_rg = rg_png.read_bytes()
        fmdx = FastMDXplora(
            system=str(project_with_analysis / "simulation" / "topology.pdb"),
            output_dir=project_with_analysis,
        )
        result = fmdx.report()

        assert result.status == "ok"
        dashboard = project_with_analysis / "report" / "dashboard.html"
        assert dashboard.is_file()
        text = dashboard.read_text(encoding="utf-8")

        assert rmsd_png.read_bytes() == original_rmsd
        assert rg_png.read_bytes() == original_rg

        assert '<aside class="sidebar">' in text
        assert "Run Progress" in text
        assert "Live Simulation" in text
        assert "Live simulation telemetry was not recorded for this run" in text
        assert "Top Metrics" in text
        assert "Recent Outputs" in text
        assert "Quick Actions" in text
        assert 'class="plot-frame"' in text
        assert 'class="plot-grid"' in text
        assert "plot-card card-md" in text
        assert "card-sm" in text
        assert "card-lg" in text
        assert "card-wide" in text
        assert "resize-handle" in text
        assert 'title="Drag to resize"' in text
        assert 'data-card-size="sm"' in text
        assert 'data-card-size="md"' in text
        assert 'data-card-size="lg"' in text
        assert 'data-card-size="wide"' in text
        assert "Reset layout" in text
        assert "data-reset-layout" in text
        assert "grid-auto-flow: dense" in text
        assert "grid-auto-rows: 8px" in text
        assert "grid-column: span var(--col-span, 1)" in text
        assert "grid-row: span var(--row-span, 20)" in text
        assert "pointerdown" in text
        assert "pointermove" in text
        assert "setPointerCapture" in text
        assert "localStorage" in text
        assert "ResizeObserver" not in text
        assert "resize: both" not in text
        assert ".panels" not in text
        assert "repeat(auto-fill, minmax(280px, 1fr))" in text
        assert "min-height: 180px" in text
        assert "overflow: hidden" in text
        assert "max-height: 100%" in text
        assert "object-fit: contain" in text
        assert "plot-card large" not in text
        assert ".plot-card.card-lg { --col-span: 2; --row-span: 38; }" in text
        assert ".plot-card.card-wide { --col-span: 2; --row-span: 24; }" in text
        assert '<span class="tag">' in text
        assert "analysis/rmsd/rmsd.png" in text
        assert "analysis/rg/rg.png" in text
        assert "Open Markdown Report" in text
        assert "Open Analysis Manifest" in text
        assert "output-list" in text
        assert "outputs-extra" in text
        assert "repeat(auto-fit, minmax(220px, 1fr))" in text
        assert "action-title" in text
        assert "action-subtitle" in text
        assert "artifact-path" in text
        assert "Std. Dev." in text
        assert "../analysis/analysis_manifest.json" in text
        assert "report.md" in text
        assert "slides.pptx" in text
        assert "project_bundle.zip" in text
        assert "dashboard.html" in text
        assert "TODO" not in text
        assert "fake" not in text.lower()
        assert "dummy" not in text.lower()
        with zipfile.ZipFile(project_with_analysis / "report" / "project_bundle.zip") as zf:
            names = set(zf.namelist())

    def test_dashboard_analysis_only_phase_aware_text(self, tmp_path: Path):
        root = tmp_path / "analysis_only_dashboard"
        analysis = root / "analysis"
        rmsd_dir = analysis / "rmsd"
        rg_dir = analysis / "rg"
        rmsd_dir.mkdir(parents=True)
        rg_dir.mkdir()
        (rmsd_dir / "rmsd.png").write_bytes(b"png")
        (rg_dir / "rg.png").write_bytes(b"png")
        (root / "manifest.json").write_text(
            json.dumps(
                {
                    "system": "1L2Y",
                    "phases": [
                        {"name": "analysis", "status": "ok"},
                        {"name": "report", "status": "ok"},
                    ],
                }
            ),
            encoding="utf-8",
        )
        (analysis / "analysis_manifest.json").write_text(
            json.dumps(
                {
                    "plan": ["rmsd", "rg"],
                    "n_frames": 10,
                    "n_atoms": 100,
                    "results": {"rmsd": {"status": "ok"}, "rg": {"status": "ok"}},
                }
            ),
            encoding="utf-8",
        )

        fmdx = FastMDXplora(system="1L2Y", output_dir=root)
        result = fmdx.report(title="Analysis-only dashboard", slides=False, bundle=False)

        assert result.status == "ok"
        text = (root / "report" / "dashboard.html").read_text(encoding="utf-8")
        assert (
            "Analysis/report workflow from existing trajectory. Setup and "
            "simulation were not run in this workflow."
        ) in text
        assert "Existing trajectory analysis" in text
        assert "<span>Setup</span><span class=\"phase-detail\">Not run</span>" in text
        assert "<span>Simulation</span><span class=\"phase-detail\">Not run</span>" in text
        assert "<span>Analysis</span><span class=\"phase-detail\">Completed</span>" in text
        assert "<span>Report</span><span class=\"phase-detail\">Completed</span>" in text
        assert 'class="plot-frame"' in text
        assert '<span class="tag">' in text
        assert "dashboard_assets" not in text
        assert "Quick Actions" in text
        assert "Recent Outputs" in text
        assert "Simulation time" not in text
        assert "Temperature" not in text
        assert "Production MD completed" not in text
        assert "../analysis/rmsd/rmsd.png" in text
        assert "../analysis/rg/rg.png" in text
        assert "project_bundle.zip" not in text

    def test_dashboard_single_artifact_sections_use_normal_plot_grid(
        self, tmp_path: Path
    ):
        root = tmp_path / "single_artifact_sections"
        artifacts = [
            "analysis/sasa/sasa.png",
            "analysis/ss/ss.png",
            "analysis/dimred/dimred_pca.png",
        ]
        png = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
            b"\x00\x00\x00\x0cIDATx\x9cc```\x00\x00\x00\x04\x00\x01"
            b"\xf6\x178U\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        for rel in artifacts:
            path = root / rel
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_bytes(png)
        (root / "manifest.json").write_text(
            json.dumps(
                {
                    "system": "1L2Y",
                    "phases": [
                        {"name": "analysis", "status": "ok"},
                        {"name": "report", "status": "ok"},
                    ],
                }
            ),
            encoding="utf-8",
        )
        (root / "analysis" / "analysis_manifest.json").write_text(
            json.dumps(
                {
                    "plan": ["sasa", "ss", "dimred"],
                    "results": {
                        "sasa": {"status": "ok"},
                        "ss": {"status": "ok"},
                        "dimred": {"status": "ok"},
                    },
                }
            ),
            encoding="utf-8",
        )

        fmdx = FastMDXplora(system="1L2Y", output_dir=root)
        result = fmdx.report(document=False, slides=False, bundle=False)

        assert result.status == "ok"
        text = (root / "report" / "dashboard.html").read_text(encoding="utf-8")
        # SASA, secondary structure, and dimensionality reduction each keep
        # their own section now, rather than being pooled because they hold
        # a single figure apiece.
        assert text.count('class="plot-grid"') == 3
        for section in ("Solvent Accessible Surface Area", "Secondary Structure",
                        "Dimensionality Reduction"):
            assert section in text
        # Each analysis anchors its own section, whatever its figure count.
        assert 'id="solvent-accessible-surface-area"' in text
        assert 'id="secondary-structure"' in text
        assert 'id="dimensionality-reduction"' in text
        assert 'id="additional-analysis"' not in text
        assert "Total SASA" in text
        assert "Secondary structure" in text
        assert "PCA" in text
        assert text.count('class="plot-card card-md"') == len(artifacts)
        assert text.count('class="resize-handle"') == len(artifacts)
        assert text.count('class="plot-frame"') == len(artifacts)
        assert "plot-card large" not in text
        assert ".plot-card.card-lg { --col-span: 2; --row-span: 38; }" in text
        assert ".plot-card.card-wide { --col-span: 2; --row-span: 24; }" in text
        assert "section-secondary-structure-section { --plot-min" not in text
        assert "section-sasa-section { --plot-min" not in text
        assert '<span class="tag">' in text
        assert "data-card-key=" in text

    def test_dashboard_multi_method_dark_assets(self, project_with_multi_method: Path):
        original_pngs = {
            path: path.read_bytes()
            for path in (
                project_with_multi_method / "analysis" / "cluster" / "cluster_kmeans.png",
                project_with_multi_method
                / "analysis"
                / "cluster"
                / "cluster_hierarchical.png",
                project_with_multi_method / "analysis" / "dimred" / "dimred_pca.png",
            )
        }
        fmdx = FastMDXplora(
            system=str(project_with_multi_method / "simulation" / "topology.pdb"),
            output_dir=project_with_multi_method,
        )
        result = fmdx.report()

        assert result.status == "ok"
        text = (
            project_with_multi_method / "report" / "dashboard.html"
        ).read_text(encoding="utf-8")
        assert "Hierarchical dendrogram" in text
        assert "analysis/cluster/cluster_hierarchical_dendrogram.png" in text
        assert '<span class="tag">' in text
        with zipfile.ZipFile(
            project_with_multi_method / "report" / "project_bundle.zip"
        ) as zf:
            names = set(zf.namelist())
        for path, data in original_pngs.items():
            assert path.read_bytes() == data

    def test_dashboard_discovers_all_analysis_image_artifacts(self, tmp_path: Path):
        root = tmp_path / "artifact_complete"
        artifact_names = [
            "analysis/sasa/total_sasa.png",
            "analysis/sasa/residue_sasa.png",
            "analysis/sasa/average_residue_sasa.png",
            "analysis/dimred/dimred_pca.png",
            "analysis/dimred/dimred_mds.png",
            "analysis/dimred/dimred_tsne.png",
            "analysis/cluster/dbscan_pop.png",
            "analysis/cluster/dbscan_traj_hist.png",
            "analysis/cluster/dbscan_traj_scatter.png",
            "analysis/cluster/dbscan_distance_matrix.png",
            "analysis/cluster/kmeans_pop.png",
            "analysis/cluster/kmeans_traj_hist.png",
            "analysis/cluster/kmeans_traj_scatter.png",
            "analysis/cluster/hierarchical_pop.png",
            "analysis/cluster/hierarchical_traj_hist.png",
            "analysis/cluster/hierarchical_traj_scatter.png",
            "analysis/cluster/hierarchical_dendrogram.png",
        ]
        png = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
            b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
            b"\x00\x00\x00\x0cIDATx\x9cc```\x00\x00\x00\x04\x00\x01"
            b"\xf6\x178U\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        for rel in artifact_names:
            path = root / rel
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_bytes(png)
        (root / "manifest.json").write_text(
            json.dumps(
                {
                    "system": "1L2Y",
                    "phases": [
                        {"name": "analysis", "status": "ok"},
                        {"name": "report", "status": "ok"},
                    ],
                }
            ),
            encoding="utf-8",
        )
        (root / "analysis" / "analysis_manifest.json").write_text(
            json.dumps(
                {
                    "plan": ["sasa", "dimred", "cluster"],
                    "results": {
                        "sasa": {"status": "ok"},
                        "dimred": {"status": "ok"},
                        "cluster": {"status": "ok"},
                    },
                }
            ),
            encoding="utf-8",
        )

        fmdx = FastMDXplora(system="1L2Y", output_dir=root)
        result = fmdx.report(document=False, slides=False, bundle=False)

        assert result.status == "ok"
        text = (root / "report" / "dashboard.html").read_text(encoding="utf-8")
        assert text.count('<div class="plot-frame">') == len(artifact_names)
        assert text.count('class="plot-card card-md"') == len(artifact_names)
        assert text.count('class="resize-handle"') == len(artifact_names)
        assert text.count('class="plot-card card-md"') == len(artifact_names)
        assert '<span class="tag">' in text
        assert "dashboard view" not in text
        assert "output-list" in text
        assert "outputs-extra" in text
        assert "Show all outputs" in text
        assert "plot-card large" not in text
        assert ".plot-card.card-lg { --col-span: 2; --row-span: 38; }" in text
        assert ".plot-card.card-wide { --col-span: 2; --row-span: 24; }" in text
        for rel in artifact_names:
            assert f"../{rel}" in text
        for label in (
            "Total SASA",
            "Per-residue SASA heatmap",
            "Average per-residue SASA",
            "PCA",
            "MDS",
            "t-SNE",
            "DBSCAN population plot",
            "DBSCAN trajectory histogram",
            "DBSCAN trajectory scatter",
            "DBSCAN distance matrix",
            "KMeans population plot",
            "KMeans trajectory histogram",
            "KMeans trajectory scatter",
            "Hierarchical population plot",
            "Hierarchical trajectory histogram",
            "Hierarchical trajectory scatter",
            "Hierarchical dendrogram",
        ):
            assert label in text
        assert "Core Metrics" not in text
        assert "SASA" in text
        assert "Dimensionality Reduction" in text
        assert "Clustering" in text


# ===========================================================================
# End-to-end with no analyses (deferred case)
# ===========================================================================
def test_report_handles_missing_analyses_gracefully(tmp_path: Path):
    """When no analysis has run, the report should still build a coherent deck."""
    fmdx = FastMDXplora(
        system="1L2Y",  # no trajectory at all
        output_dir=tmp_path,
    )
    result = fmdx.report()
    assert result.status == "ok"
    # The deck should exist; image-slide count is zero because no analyses ran
    from pptx import Presentation

    prs = Presentation(str(tmp_path / "report" / "slides.pptx"))
    assert len(prs.slides) >= 1


def test_deferred_analysis_message_is_current_and_actionable(tmp_path: Path):
    fmdx = FastMDXplora(system="1L2Y", output_dir=tmp_path)
    result = fmdx.analyze()

    assert result.status == "ok"
    manifest = json.loads(
        (tmp_path / "analysis" / "analysis_manifest.json").read_text(encoding="utf-8")
    )
    assert manifest["status"] == "deferred"
    assert "v0.2" not in manifest["note"]
    assert "Run the simulation phase first" in manifest["note"]


# ===========================================================================
# Robustness for report-only mode and generated artifacts
# ===========================================================================
def test_report_only_cli_escapes_weird_markdown_text(tmp_path: Path):
    """Report-only mode should not let titles/messages reshape Markdown."""
    root = tmp_path / "existing run with spaces"
    analysis = root / "analysis"
    analysis.mkdir(parents=True)
    weird_system = "sys|`tick` [brackets]\n# injected\tpath"
    weird_title = "Study | `ticks` [link](x)\n# injected"
    weird_author = "Author | # name\nnext"
    weird_message = "bad | `tick`\n# injected failure"

    (root / "manifest.json").write_text(
        json.dumps({"system": weird_system}),
        encoding="utf-8",
    )
    (analysis / "analysis_manifest.json").write_text(
        json.dumps(
            {
                "plan": ["rmsd"],
                "results": {"rmsd": {"status": "error", "message": weird_message}},
            }
        ),
        encoding="utf-8",
    )

    rc = cli_main(
        [
            "report",
            "--output",
            str(root),
            "--title",
            weird_title,
            "--author",
            weird_author,
            "--no-slides",
            "--no-bundle",
        ]
    )

    assert rc == 0
    text = (root / "report" / "report.md").read_text(encoding="utf-8")
    first_line = text.splitlines()[0]
    assert first_line.startswith("# Study")
    assert "\\|" in first_line
    assert "\\`ticks\\`" in first_line
    assert "\n# injected" not in text
    assert "bad \\| \\`tick\\` \\# injected failure" in text


def test_report_phase_manifest_artifacts_exist(tmp_path: Path):
    fmdx = FastMDXplora(system="sys|odd", output_dir=tmp_path / "run")
    results = fmdx.explore(
        include=["report"],
        options={"report": {"title": "Odd | title", "slides": False, "bundle": False}},
    )

    report_phase = results[0].phase("report")
    assert report_phase is not None
    manifest = json.loads((fmdx.output_dir / "manifest.json").read_text(encoding="utf-8"))
    report_record = next(p for p in manifest["phases"] if p["name"] == "report")
    # A PDF is written where WeasyPrint is available and left out where it is
    # not, so the list is checked for what must be there rather than for an
    # exact sequence that depends on what somebody installed. Where it is left
    # out, not_produced.json says so -- the two are alternatives, and exactly
    # one of them appears.
    written = report_record["artifacts"]
    assert "report.md" in written and "dashboard.html" in written
    assert set(written) <= {"report.md", "report.pdf", "dashboard.html",
                            "not_produced.json"}
    assert ("report.pdf" in written) != ("not_produced.json" in written), (
        "a PDF was requested, so either it was written or the reason it was "
        "not is recorded beside the outputs")
    for artifact in report_record["artifacts"]:
        assert (Path(report_record["output_dir"]) / artifact).is_file()


def test_report_bundle_excludes_cache_and_temp_files(tmp_path: Path):
    root = tmp_path / "run"
    (root / "analysis").mkdir(parents=True)
    (root / "__pycache__").mkdir()
    (root / ".cache").mkdir()
    (root / "__pycache__" / "module.cpython-310.pyc").write_bytes(b"cached")
    (root / ".cache" / "plot.tmp").write_text("cache", encoding="utf-8")
    (root / "scratch.tmp").write_text("temp", encoding="utf-8")
    (root / "keep.dat").write_text("artifact", encoding="utf-8")

    fmdx = FastMDXplora(system="1L2Y", output_dir=root)
    result = fmdx.report(document=False, slides=False, bundle=True)

    assert result.status == "ok"
    bundle = root / "report" / "project_bundle.zip"
    with zipfile.ZipFile(bundle) as zf:
        names = set(zf.namelist())
    assert "keep.dat" in names
    assert "report/dashboard.html" in names
    assert "scratch.tmp" not in names
    assert "__pycache__/module.cpython-310.pyc" not in names
    assert ".cache/plot.tmp" not in names
    assert "report/project_bundle.zip" not in names


def test_powerpoint_handles_weird_title_and_system(tmp_path: Path):
    from pptx import Presentation

    title = "T" * 800 + " | `tick`\n# heading"
    system = "system with spaces | [id]\nnext"
    fmdx = FastMDXplora(system=system, output_dir=tmp_path / "run")
    result = fmdx.report(title=title, bundle=False)

    assert result.status == "ok"
    pptx = fmdx.output_dir / "report" / "slides.pptx"
    assert pptx.is_file()
    prs = Presentation(str(pptx))
    assert len(prs.slides) >= 1
    assert "\n# heading" not in prs.slides[0].shapes.title.text


def test_analysis_report_only_wording_uses_existing_trajectory(tmp_path: Path):
    from pptx import Presentation

    root = tmp_path / "analysis_only"
    analysis = root / "analysis"
    analysis.mkdir(parents=True)
    (root / "manifest.json").write_text(
        json.dumps(
            {
                "system": "1L2Y",
                "phases": [
                    {"name": "analysis", "status": "ok"},
                    {"name": "report", "status": "ok"},
                ],
            }
        ),
        encoding="utf-8",
    )
    (analysis / "analysis_manifest.json").write_text(
        json.dumps(
            {
                "plan": ["rmsd"],
                "n_frames": 10,
                "n_residues": 20,
                "results": {"rmsd": {"status": "ok"}},
            }
        ),
        encoding="utf-8",
    )
    rmsd_dir = analysis / "rmsd"
    rmsd_dir.mkdir()
    (rmsd_dir / "options.json").write_text(
        json.dumps({"selection": "name CA", "options": {"align": True}}),
        encoding="utf-8",
    )

    fmdx = FastMDXplora(system="1L2Y", output_dir=root)
    result = fmdx.report(title="Analysis-only Trp-cage report", bundle=False)

    assert result.status == "ok"
    report = (root / "report" / "report.md").read_text(encoding="utf-8")
    assert "end-to-end molecular dynamics study" not in report
    assert "Simulation parameters were not recorded for this run" not in report
    assert "This report was generated from an existing trajectory" in report
    assert "Setup and simulation were not run in this workflow" in report
    assert "Simulation was not run in this workflow" in report

    outline = (root / "report" / "slides_outline.md").read_text(encoding="utf-8")
    assert "Analysis/report workflow from an existing trajectory" in outline
    assert "Setup and simulation were not run in this workflow" in outline
    assert "## 2. Setup" not in outline
    assert "## 3. Simulation" not in outline

    prs = Presentation(str(root / "report" / "slides.pptx"))
    all_text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Analysis/report workflow from an existing trajectory" in all_text
    assert "Setup and simulation were not run in this workflow" in all_text


def test_full_pipeline_report_keeps_end_to_end_wording(tmp_path: Path):
    root = tmp_path / "full_pipeline"
    (root / "setup").mkdir(parents=True)
    (root / "simulation").mkdir()
    (root / "analysis").mkdir()
    (root / "manifest.json").write_text(
        json.dumps(
            {
                "system": "1L2Y",
                "phases": [
                    {"name": "setup", "status": "ok"},
                    {"name": "simulation", "status": "ok"},
                    {"name": "analysis", "status": "ok"},
                    {"name": "report", "status": "ok"},
                ],
            }
        ),
        encoding="utf-8",
    )
    (root / "setup" / "setup_parameters.json").write_text(
        json.dumps({"parameters": {"ph": 7.0}}),
        encoding="utf-8",
    )
    (root / "simulation" / "simulation_parameters.json").write_text(
        json.dumps({"parameters": {"duration_ns": 0.001}}),
        encoding="utf-8",
    )
    (root / "analysis" / "analysis_manifest.json").write_text(
        json.dumps({"plan": [], "results": {}}),
        encoding="utf-8",
    )

    fmdx = FastMDXplora(system="1L2Y", output_dir=root)
    result = fmdx.report(title="Full pipeline report", slides=False, bundle=False)

    assert result.status == "ok"
    report = (root / "report" / "report.md").read_text(encoding="utf-8")
    # The summary used to open with "end-to-end molecular dynamics study" for
    # every full run. It now says what the study was, so the distinction this
    # test is about -- a full pipeline against analysis from a supplied
    # trajectory -- is checked by what it does not claim.
    assert "1 ps" in report, "the summary states the run it is about"
    assert "Production MD was performed" in report
    assert "Setup and simulation were not run in this workflow" not in report


def test_region_highlight_report_from_existing_rmsf_outputs(tmp_path: Path):
    from pptx import Presentation

    root = tmp_path / "region_run"
    rmsf_dir = root / "analysis" / "rmsf"
    rmsf_dir.mkdir(parents=True)
    np.savetxt(
        rmsf_dir / "rmsf.dat",
        np.column_stack([np.arange(1, 9), np.linspace(0.02, 0.12, 8)]),
    )
    (root / "manifest.json").write_text(
        json.dumps(
            {
                "system": "example.pdb",
                "phases": [
                    {"name": "analysis", "status": "ok"},
                    {"name": "report", "status": "ok"},
                ],
            }
        ),
        encoding="utf-8",
    )
    (root / "analysis" / "analysis_manifest.json").write_text(
        json.dumps(
            {
                "plan": ["rmsf"],
                "n_frames": 5,
                "n_residues": 8,
                "topology_input": None,
                "results": {"rmsf": {"status": "ok"}},
            }
        ),
        encoding="utf-8",
    )

    fmdx = FastMDXplora(system="example.pdb", output_dir=root)
    result = fmdx.report(
        bundle=False,
        region_highlights=[
            {"label": "example region 1", "start": 2, "end": 4, "color": "#4E79A7"},
            {"label": "example region 2", "start": 6, "end": 7, "color": "#F28E2B"},
        ],
    )

    assert result.status == "ok"
    assert (root / "analysis" / "rmsf" / "rmsf_region_highlights.png").is_file()
    assert (root / "report" / "region_highlight_summary.png").is_file()
    manifest = json.loads(
        (root / "report" / "region_highlight_manifest.json").read_text(encoding="utf-8")
    )
    assert manifest["status"] == "ok"
    assert manifest["skipped"][0]["artifact"] == "structure_region_highlights.png"
    reason = manifest["skipped"][0]["reason"]
    assert "PyMOL" in reason or "topology" in reason

    report = (root / "report" / "report.md").read_text(encoding="utf-8")
    assert "Region Highlight Figure" in report
    assert "region_highlight_summary.png" in report
    outline = (root / "report" / "slides_outline.md").read_text(encoding="utf-8")
    assert "region_highlight_summary.png" in outline
    prs = Presentation(str(root / "report" / "slides.pptx"))
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title is not None]
    assert "Region highlights" in titles


def test_region_highlight_invalid_range_records_error(tmp_path: Path):
    root = tmp_path / "bad_region"
    rmsf_dir = root / "analysis" / "rmsf"
    rmsf_dir.mkdir(parents=True)
    np.savetxt(rmsf_dir / "rmsf.dat", np.column_stack([np.arange(1, 5), np.ones(4)]))
    (root / "analysis" / "analysis_manifest.json").write_text(
        json.dumps({"plan": ["rmsf"], "results": {"rmsf": {"status": "ok"}}}),
        encoding="utf-8",
    )

    fmdx = FastMDXplora(system="example.pdb", output_dir=root)
    result = fmdx.report(
        document=False,
        slides=False,
        bundle=False,
        region_highlights=[{"label": "bad", "start": 0, "end": 2}],
    )

    assert result.status == "ok"
    manifest = json.loads(
        (root / "report" / "region_highlight_manifest.json").read_text(encoding="utf-8")
    )
    assert manifest["status"] == "error"
    assert "start must be >= 1" in manifest["error"]
    assert not (root / "analysis" / "rmsf" / "rmsf_region_highlights.png").exists()


def test_report_without_region_highlights_has_no_region_artifacts(
    project_with_analysis: Path,
):
    fmdx = FastMDXplora(
        system=str(project_with_analysis / "simulation" / "topology.pdb"),
        output_dir=project_with_analysis,
    )
    fmdx.report(bundle=False)

    assert not (project_with_analysis / "report" / "region_highlight_summary.png").exists()
    text = (project_with_analysis / "report" / "report.md").read_text(encoding="utf-8")
    assert "Region Highlight Figure" not in text


def test_region_highlight_structure_panel_when_pymol_available(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
):
    import subprocess

    from fastmdxplora.report import region_highlights as rh

    def fake_run(cmd, check, text, capture_output, timeout):
        import matplotlib.pyplot as plt

        script_path = Path(cmd[-1])
        script = script_path.read_text(encoding="utf-8")
        out_line = next(line for line in script.splitlines() if line.startswith("png "))
        out_path = Path(out_line.split(" ", 1)[1].split(",", 1)[0])
        plt.imsave(out_path, np.ones((4, 4, 3)))
        return subprocess.CompletedProcess(cmd, 0, "", "")

    monkeypatch.setattr(
        rh,
        "detect_pymol_renderer",
        lambda: rh.PymolRenderer(kind="command", command=("pymol", "-cq")),
    )
    monkeypatch.setattr(rh.subprocess, "run", fake_run)

    root = _make_traj_files(tmp_path / "region_structure", n_residues=8, n_frames=10)
    fmdx = FastMDXplora(
        system=str(root / "simulation" / "topology.pdb"),
        output_dir=root,
    )
    fmdx.explore(
        include=["analysis", "report"],
        options={
            "analysis": {"include": ["rmsf"]},
            "report": {
                "bundle": False,
                "region_highlights": [
                    {"label": "example region", "start": 2, "end": 5}
                ],
            },
        },
    )

    assert (root / "report" / "structure_region_highlights.png").is_file()
    assert (root / "report" / "structure_region_highlights.pml").is_file()
    assert (root / "report" / "region_highlight_summary.png").is_file()
    manifest = json.loads(
        (root / "report" / "region_highlight_manifest.json").read_text(encoding="utf-8")
    )
    assert manifest["renderer"] == "PyMOL"
    assert "report/structure_region_highlights.png" in manifest["artifacts"]
    assert "report/structure_region_highlights.pml" in manifest["artifacts"]
    assert manifest["skipped"] == []


def test_pymol_script_generation_uses_cartoon_and_regions(tmp_path: Path):
    from fastmdxplora.report.region_highlights import (
        RegionHighlight,
        build_pymol_script,
    )

    script = build_pymol_script(
        topology_path=tmp_path / "topology.pdb",
        output_path=tmp_path / "structure.png",
        regions=[
            RegionHighlight("example A", 2, 5, "#4E79A7"),
            RegionHighlight("example B", 7, 9, "orange"),
        ],
    )

    assert "show cartoon, prot" in script
    assert "color gray70, prot" in script
    assert "resi 2-5" in script
    assert "resi 7-9" in script
    assert "ray 1800, 1200" in script
    assert "png " in script


class TestTheReportWritesAMethodsSection:
    """A journal asks for a specific list of things, and the list is published:
    JCIM's Guidelines for Reporting Molecular Dynamics Simulations (Soares et
    al. 2023) and Communications Biology's reproducibility checklist (2023).

    Every value was already recorded here. What was missing was the assembly --
    the report listed each setting as a bullet, which is what the software
    knows rather than what a reader needs. Nobody types this paragraph
    correctly from memory, which is why the drMD authors described difficulty
    reproducing a published protocol not for technical reasons but because of
    ambiguity in how it was presented.
    """

    @staticmethod
    def _run():
        # The shape a real run writes: `parameters` holds what was asked for,
        # and what the run resolved sits beside it. A flat fixture would let
        # this pass while the real thing failed, which is how the first
        # version shipped saying the force field was "amber-openff".
        setup = {
            "input": {"system": "181L", "form": "pdb_id"},
            "parameters": {
                "ph": 7.4, "heterogens": "auto", "forcefield": "amber-openff",
                # Written by setup when it prepares a ligand; a manifest
                # without it is a run that had none.
                "ligand": ["ligands/BNZ.sdf"],
                "ligand_name": "BNZ", "ligand_net_charge": 0,
                "solvent_padding_nm": 0.8, "box_shape": "cube",
                "ion_concentration_M": 0.15, "nonbonded_method": "PME",
                "nonbonded_cutoff_nm": 1.0, "constraints": "HBonds",
                "rigid_water": True,
            },
            "resolved_forcefield": {
                "xmls": ["amber14-all.xml", "amber14/tip3p.xml"],
                "water_model": "tip3p",
                "small_molecule_forcefield": "openff-2.2.1",
            },
            "n_atoms_solvated": 34020,
        }
        sim = {
            "parameters": {
                "timestep_fs": 2.0, "integrator": "langevin_middle",
                "temperature_K": 300.0, "friction_per_ps": 1.0,
                "pressure_bar": 1.0, "barostat_frequency": 25,
                "nvt_steps": 200, "npt_steps": 200, "production_steps": 2000,
                "trajectory_interval_steps": 20, "precision": "mixed",
                "minimize": True, "random_seed": 7,
            },
            "platform_used": "CPU",
        }
        return setup, sim

    def test_it_states_what_the_checklists_require(self) -> None:
        from pathlib import Path

        from fastmdxplora.report.methods import methods_paragraphs

        setup, sim = self._run()
        text = methods_paragraphs(Path("."), setup, sim)
        for required in ("181L", "pH 7.4", "amber14", "tip3p", "openff-2.2.1",
                         "PME", "HBonds", "Langevin", "300.0 K",
                         "1.0 bar", "34,020"):
            assert required in text, f"the methods section omits {required}"

    def test_steps_are_given_as_time(self) -> None:
        """Nobody reads "2000 steps". The reader needs to know it was 4 ps."""
        from pathlib import Path

        from fastmdxplora.report.methods import methods_paragraphs

        setup, sim = self._run()
        text = methods_paragraphs(Path("."), setup, sim)
        assert "4 ps" in text
        assert "2000 steps" not in text

    def test_a_missing_seed_is_reported_as_missing(self) -> None:
        """Its absence is what makes a run irreproducible, so it is as
        reportable as its value."""
        from pathlib import Path

        from fastmdxplora.report.methods import methods_paragraphs

        setup, sim = self._run()
        sim["parameters"]["random_seed"] = None
        text = methods_paragraphs(Path("."), setup, sim)
        assert "No random seed was fixed" in text

    def test_what_was_not_recorded_is_named(self) -> None:
        """A methods section quietly stating a default nobody chose is worse
        than one with a visible gap."""
        from pathlib import Path

        from fastmdxplora.report.methods import (
            methods_paragraphs,
            missing_from_methods,
        )

        setup, sim = self._run()
        del setup["resolved_forcefield"]["water_model"]
        assert "water model" in missing_from_methods(setup, sim)
        text = methods_paragraphs(Path("."), setup, sim)
        assert "Not recorded" in text
        assert "water model" in text

    def test_nothing_is_invented(self) -> None:
        """An empty run produces no claims about it."""
        from pathlib import Path

        from fastmdxplora.report.methods import methods_paragraphs

        text = methods_paragraphs(Path("."), {}, {})
        assert "PME" not in text and "300" not in text
        assert "Not recorded" in text

    def test_the_software_is_named_with_its_version(self) -> None:
        """Naming a tool without its version names a moving target: the
        defaults change and a reader gets different numbers with no way to
        know why."""
        from fastmdxplora.report.document import _software_versions

        found = _software_versions()
        assert found, "nothing reported a version"
        assert all(version for version in found.values())

    def test_the_report_puts_the_prose_before_the_settings(self) -> None:
        import inspect

        from fastmdxplora.report import document

        source = inspect.getsource(document._methods_section)
        prose = source.index("methods_paragraphs")
        settings = source.index("### System preparation")
        assert prose < settings, (
            "the settings list should follow the paragraph, not replace it"
        )


class TestTheMethodsSectionReadsTheRealManifests:
    """The first version read only ``parameters`` and produced a methods
    section saying the coordinates came from "the input structure", the force
    field was "amber-openff", and the water model was not recorded.

    None of that was true. The manifests nest: the system sits under
    ``input``, and what the run resolved -- the actual XML files, the water
    model, the small-molecule force field -- sits beside ``parameters`` rather
    than inside it. Guessing at key names produced a paragraph that read well
    and said the wrong thing, which is worse than one that reads badly.
    """

    @staticmethod
    def _manifests():
        """The shape a real run writes, not a convenient flat one."""
        setup = {
            "input": {"system": "181L", "form": "pdb_id"},
            "parameters": {
                "ph": 7.4, "heterogens": "auto", "forcefield": "amber-openff",
                "solvent_padding_nm": 0.8, "box_shape": "cube",
                "ion_concentration_M": 0.15, "nonbonded_method": "PME",
                "nonbonded_cutoff_nm": 1.0, "constraints": "HBonds",
                "rigid_water": True, "ligand": ["ligands/BNZ.sdf"],
                "ligand_name": "BNZ", "ligand_net_charge": 0,
            },
            "resolved_forcefield": {
                "source": "named", "name": "amber-openff",
                "xmls": ["amber14-all.xml", "amber14/tip3p.xml"],
                "water_model": "tip3p",
                "small_molecule_forcefield": "openff-2.2.1",
            },
            "n_atoms_solvated": 35012,
        }
        sim = {
            "parameters": {
                "timestep_fs": 2.0, "integrator": "langevin_middle",
                "temperature_K": 300.0, "friction_per_ps": 1.0,
                "pressure_bar": 1.0, "barostat_frequency": 25,
                "nvt_steps": 200, "npt_steps": 200, "production_steps": 2000,
                "trajectory_interval_steps": 20, "platform": "auto",
                "precision": "mixed", "minimize": True, "random_seed": None,
            },
            "platform_used": "CPU", "n_production_frames": 100,
        }
        return setup, sim

    def _text(self):
        from pathlib import Path

        from fastmdxplora.report.methods import methods_paragraphs

        setup, sim = self._manifests()
        return methods_paragraphs(Path("."), setup, sim,
                                  system_name=setup["input"]["system"])

    def test_the_system_is_named(self) -> None:
        text = self._text()
        assert "181L" in text
        assert "the input structure" not in text

    def test_the_force_field_is_the_one_used_not_our_label(self) -> None:
        """A reader cannot look up "amber-openff"; they can look up
        amber14-all.xml."""
        text = self._text()
        assert "amber14-all.xml" in text
        assert "tip3p" in text

    def test_the_ligand_and_its_parameters_are_stated(self) -> None:
        text = self._text()
        assert "BNZ" in text
        assert "openff-2.2.1" in text

    def test_the_system_size_is_stated(self) -> None:
        text = self._text()
        assert "35,012" in text

    def test_the_platform_is_what_ran_not_what_was_asked_for(self) -> None:
        """A methods section saying a run used the "auto" platform says
        nothing."""
        text = self._text()
        assert "CPU platform" in text
        assert "auto platform" not in text

    def test_the_integrator_is_named_as_the_literature_names_it(self) -> None:
        text = self._text()
        assert "Langevin middle-scheme" in text
        assert "langevin_middle" not in text

    def test_the_barostat_is_stated(self) -> None:
        text = self._text()
        assert "1.0 bar" in text and "every 25 steps" in text

    def test_nothing_present_is_reported_as_missing(self) -> None:
        """The water model was recorded, and saying otherwise sends somebody
        looking for a number that is already there."""
        text = self._text()
        assert "Not recorded" not in text

    def test_the_report_passes_the_whole_manifest(self) -> None:
        import inspect

        from fastmdxplora.report import document

        source = inspect.getsource(document._methods_section)
        call = source[source.index("methods_paragraphs("):][:300]
        assert "project_root, setup, sim" in call, (
            "passing only `parameters` loses the system and the resolution"
        )


class TestWhatTheRunResolvedIsRecorded:
    """Twice the methods section had to report something as unrecorded that
    the run plainly knew.

    The solvated atom count was printed to the terminal and discarded. The
    pressure the barostat ran at was resolved inside the runner -- it can be
    given in bar or atmospheres, and unset means one bar -- and never left it.
    Both are on the checklist a journal applies, so a run of an NPT simulation
    could not say what pressure it held.
    """

    def test_preparation_reports_the_system_size(self) -> None:
        import inspect

        from fastmdxplora.setup import prepare

        returned = inspect.getsource(prepare.prepare_system).rsplit("return", 1)[-1]
        assert "n_atoms_solvated" in returned

    def test_the_setup_manifest_records_it(self) -> None:
        import inspect

        from fastmdxplora.setup import pipeline

        source = inspect.getsource(pipeline._write_manifest)
        assert '"n_atoms_solvated"' in source

    def test_the_runner_reports_the_pressure_it_used(self) -> None:
        import inspect

        from fastmdxplora.simulation import runner

        source = inspect.getsource(runner)
        assert "pressure_bar_used" in source
        assert "pressure_bar_used=resolved_pressure_bar" in source, (
            "it should report what it resolved, not what it was passed"
        )

    def test_the_simulation_manifest_records_it(self) -> None:
        import inspect

        from fastmdxplora.simulation import pipeline

        assert '"pressure_bar_used"' in inspect.getsource(pipeline._write_manifest)

    def test_the_methods_section_prefers_what_ran(self) -> None:
        """A run asked for nothing and held one bar. Reporting the request
        would say nothing; reporting the resolution says what happened."""
        from pathlib import Path

        from fastmdxplora.report.methods import methods_paragraphs

        sim = {
            "parameters": {"pressure_bar": None, "timestep_fs": 2.0,
                           "npt_steps": 200, "integrator": "langevin_middle"},
            "pressure_bar_used": 1.0,
        }
        text = methods_paragraphs(Path("."), {}, sim)
        assert "1.0 bar" in text

    def test_and_says_so_when_neither_is_known(self) -> None:
        from pathlib import Path

        from fastmdxplora.report.methods import missing_from_methods

        sim = {"parameters": {"timestep_fs": 2.0, "integrator": "verlet"}}
        gaps = missing_from_methods({}, sim)
        assert any("pressure" in gap for gap in gaps)


class TestConvergenceSaysWhatARunCanSupport:
    """A hundred frames look like a hundred measurements.

    Consecutive frames of a molecular dynamics run are almost the same
    structure, so an error bar computed over them is too small by the square
    root of the correlation time -- a factor of three or four is ordinary.
    This is the same mistake fixed for interaction occupancies, in the
    analyses everybody reads first.
    """

    def test_uncorrelated_data_is_counted_in_full(self) -> None:
        import numpy as np

        from fastmdxplora.report.convergence import autocorrelation_time

        rng = np.random.RandomState(0)
        assert autocorrelation_time(rng.normal(size=1000)) < 2.0

    def test_correlated_data_is_not(self) -> None:
        import numpy as np

        from fastmdxplora.report.convergence import autocorrelation_time

        rng = np.random.RandomState(0)
        smoothed = np.convolve(rng.normal(size=1200), np.ones(40) / 40,
                               mode="valid")
        assert autocorrelation_time(smoothed) > 10.0

    def test_the_error_bar_counts_independent_samples(self) -> None:
        """Dividing by the frame count would understate it several-fold."""
        import numpy as np

        from fastmdxplora.report.convergence import assess_series

        rng = np.random.RandomState(0)
        smoothed = np.convolve(rng.normal(size=1200), np.ones(40) / 40,
                               mode="valid")[:400]
        found = assess_series("correlated", smoothed)
        naive = found.spread / np.sqrt(found.n_frames)
        assert found.standard_error > naive * 2, (
            "the honest error bar should be markedly larger than the naive one"
        )

    def test_a_run_too_short_to_measure_correlation_says_so(self) -> None:
        """The estimate truncates at half the series, so a short run cannot
        see a memory longer than that -- and would report the independence it
        failed to rule out, which is the wrong direction for a convergence
        report to err in.
        """
        import numpy as np

        from fastmdxplora.report.convergence import assess_series

        rng = np.random.RandomState(1)
        slow = np.convolve(rng.normal(size=400), np.ones(60) / 60,
                           mode="valid")[:60]
        found = assess_series("slow", slow)
        assert not found.correlation_is_measurable
        assert not found.is_sampled_enough, (
            "independence that could not be measured must not count as found"
        )

    def test_a_drifting_observable_is_reported_as_unsettled(self) -> None:
        import numpy as np

        from fastmdxplora.report.convergence import assess_series

        rng = np.random.RandomState(2)
        climbing = np.linspace(0, 5, 300) + rng.normal(scale=0.1, size=300)
        assert not assess_series("climbing", climbing).has_settled

    def test_a_settled_observable_is_not(self) -> None:
        import numpy as np

        from fastmdxplora.report.convergence import assess_series

        rng = np.random.RandomState(3)
        assert assess_series("flat", rng.normal(size=300)).has_settled

    def test_a_run_that_cannot_be_interpreted_says_so_plainly(self) -> None:
        """Returning a number for everything would launder a short run into
        an apparently validated one, which is worse than no report: a reader
        would have less reason to doubt than they started with."""
        import numpy as np

        from fastmdxplora.report.convergence import assess_run

        rng = np.random.RandomState(4)
        climbing = np.linspace(0, 5, 120) + rng.normal(scale=0.1, size=120)
        out = assess_run({"rmsd": climbing})
        assert out["interpretable"] is False
        assert out["findings"]

    def test_a_thermostat_off_its_target_is_reported(self) -> None:
        import numpy as np

        from fastmdxplora.report.convergence import assess_run

        rng = np.random.RandomState(5)
        out = assess_run({"temperature": 340 + rng.normal(scale=2, size=500)},
                         target_temperature_K=300.0)
        assert any("thermostat" in finding for finding in out["findings"])

    def test_a_thermostat_holding_its_target_is_not(self) -> None:
        import numpy as np

        from fastmdxplora.report.convergence import assess_run

        rng = np.random.RandomState(6)
        out = assess_run({"temperature": 300 + rng.normal(scale=2, size=500)},
                         target_temperature_K=300.0)
        assert not any("thermostat" in finding for finding in out["findings"])


class TestTheReportSaysHowMuchTheRunSupports:
    """The convergence assessment reaches the report, after the results,
    because it is about them: every mean and error bar above rests on how many
    independent observations the run holds, and that is usually far fewer than
    the frame count suggests.
    """

    @staticmethod
    def _a_run(tmp_path, n_frames=400):
        """Analyses written where a real run writes them."""
        import mdtraj as md
        import numpy as np

        import fastmdxplora.analysis  # noqa: F401
        from fastmdxplora.analysis.orchestrator import get_analysis_class

        top = md.Topology()
        chain = top.add_chain()
        previous_c = None
        for index in range(6):
            residue = top.add_residue("ALA", chain, resSeq=index + 1)
            n = top.add_atom("N", md.element.nitrogen, residue)
            ca = top.add_atom("CA", md.element.carbon, residue)
            c = top.add_atom("C", md.element.carbon, residue)
            o = top.add_atom("O", md.element.oxygen, residue)
            for first, second in ((n, ca), (ca, c), (c, o)):
                top.add_bond(first, second)
            if previous_c is not None:
                top.add_bond(previous_c, n)
            previous_c = c

        rng = np.random.RandomState(1)
        xyz = rng.normal(scale=0.2, size=(n_frames, top.n_atoms, 3))
        traj = md.Trajectory(xyz=xyz.astype(np.float32), topology=top)
        for name in ("rmsd", "rg"):
            get_analysis_class(name)(
                output_dir=str(tmp_path / "analysis")).run(traj)
        return tmp_path

    def test_the_section_appears(self, tmp_path) -> None:
        from fastmdxplora.report.document import _convergence_section

        text = _convergence_section(self._a_run(tmp_path))
        assert "## Convergence" in text
        assert "rmsd" in text and "rg" in text

    def test_it_reports_independent_samples_not_frames(self, tmp_path) -> None:
        from fastmdxplora.report.document import _convergence_section

        text = _convergence_section(self._a_run(tmp_path))
        assert "independent" in text
        assert "A frame is not an observation" in text

    def test_it_reads_a_bare_array_as_well_as_a_table(self, tmp_path) -> None:
        """An analysis returning an array writes plain numbers; one returning
        a frame writes a header. Both shapes are read rather than one
        assumed."""
        from fastmdxplora.report.document import _last_numeric_column

        plain = tmp_path / "plain.dat"
        plain.write_text("1.0\n2.0\n3.0\n", encoding="utf-8")
        assert _last_numeric_column(plain) == [1.0, 2.0, 3.0]

        table = tmp_path / "table.dat"
        table.write_text("frame,value\n0,1.0\n1,2.0\n", encoding="utf-8")
        assert _last_numeric_column(table) == [1.0, 2.0]

    def test_a_run_with_no_data_produces_no_section(self, tmp_path) -> None:
        """Rather than a heading over an empty table."""
        from fastmdxplora.report.document import _convergence_section

        assert _convergence_section(tmp_path) == ""

    def test_it_sits_after_the_results(self) -> None:
        import inspect

        from fastmdxplora.report import document

        source = inspect.getsource(document)
        assert (source.index("_results_section(project_root))")
                < source.index("_convergence_section(project_root)"))


class TestTheReportIsAlsoAPDF:
    """A Markdown file is not a document somebody sends on.

    It renders differently in every viewer, it cannot be printed with the
    figures where the text put them, and attaching one invites the question of
    what to open it with. The package docstring had claimed a PDF for some
    time; nothing produced one.
    """

    @staticmethod
    def _report(tmp_path):
        path = tmp_path / "report.md"
        path.write_text(
            "# A report\n\n## Methods\n\nThe protein was described by the "
            "amber14-all.xml force field.\n\n## Convergence\n\n"
            "| measure | frames | independent |\n|---|---|---|\n"
            "| rmsd | 4,999 | 20.3 |\n",
            encoding="utf-8")
        return path

    def test_it_renders(self, tmp_path) -> None:
        import pytest

        pytest.importorskip("weasyprint", reason="requires the [pdf] extra")
        pytest.importorskip("markdown", reason="requires the [pdf] extra")

        from fastmdxplora.report.pdf import render_pdf

        written = render_pdf(self._report(tmp_path))
        assert written.suffix == ".pdf"
        assert written.read_bytes().startswith(b"%PDF")

    def test_it_is_written_beside_the_markdown(self, tmp_path) -> None:
        import pytest

        pytest.importorskip("weasyprint", reason="requires the [pdf] extra")
        pytest.importorskip("markdown", reason="requires the [pdf] extra")

        from fastmdxplora.report.pdf import render_pdf

        source = self._report(tmp_path)
        assert render_pdf(source).parent == source.parent

    def test_a_missing_report_says_so(self, tmp_path) -> None:
        import pytest

        from fastmdxplora.report.pdf import PdfUnavailable, render_pdf

        with pytest.raises(PdfUnavailable, match="No report to convert"):
            render_pdf(tmp_path / "nothing.md")

    def test_it_does_not_fail_the_run_when_it_cannot(self, tmp_path) -> None:
        """A phase that produced a document, a dashboard, slides and a bundle
        should not be marked failed because a fifth format needed a library
        nobody installed."""
        from fastmdxplora.report.pdf import try_render_pdf

        written, reason = try_render_pdf(tmp_path / "nothing.md")
        assert written is None
        assert reason and "No report to convert" in reason

    def test_the_reason_names_what_to_install(self) -> None:
        """Somebody told a shared library is missing has been told the
        symptom, not the remedy."""
        import inspect

        from fastmdxplora.report import pdf

        source = inspect.getsource(pdf)
        assert "fastmdxplora[pdf]" in source
        assert "conda-forge" in source
        assert "Pango" in source and "Cairo" in source

    def test_it_converts_the_document_not_the_dashboard(self) -> None:
        """A PDF of the browser dashboard would be a picture of an interface.
        The report is a document, and a document is what a PDF should be."""
        import inspect

        from fastmdxplora.report import pdf

        source = inspect.getsource(pdf)
        assert "report.md" in source
        assert "dashboard.html" not in source

    def test_the_setting_is_declared_like_the_others(self) -> None:
        from fastmdxplora.config.schema import PHASE_SCHEMAS

        field = PHASE_SCHEMAS["report"].get("pdf")
        assert field is not None and field.default is True
        assert field.help and "WeasyPrint" in field.help

    def test_the_phase_asks_for_it(self) -> None:
        import inspect

        from fastmdxplora.report import run

        source = inspect.getsource(run)
        assert "try_render_pdf" in source
        assert 'params.get("pdf"' in source


class TestAShortSeriesSaysWhatItCannotJudge:
    """A real run reported "settled: yes" for a two-point series, beside "too
    few independent samples" about the same numbers.

    Two points cannot show a trend. Reporting one as settled is a claim from
    no evidence, and it is the same defect as reporting independence that was
    never measured -- which this module was written to avoid.
    """

    def test_two_points_cannot_say_whether_it_settled(self) -> None:
        from fastmdxplora.report.convergence import assess_series

        found = assess_series("potential_energy", [-490000.0, -476830.0])
        assert found.has_settled is None
        assert not found.drift_is_measurable

    def test_a_long_series_still_can(self) -> None:
        import numpy as np

        from fastmdxplora.report.convergence import assess_series

        rng = np.random.RandomState(0)
        assert assess_series("t", 300 + rng.normal(scale=2, size=400)).has_settled

    def test_energy_drift_needs_more_than_two_samples(self) -> None:
        """It reported ninety-four kJ/mol per ns per atom from two samples of
        a run that had just been minimized. The range of two numbers is not a
        drift rate."""
        from fastmdxplora.report.convergence import assess_run

        out = assess_run({"potential_energy": [-490000.0, -476830.0]},
                         duration_ns=0.004, n_atoms=34979)
        assert not any("kJ/mol per ns per atom" in f for f in out["findings"])

    def test_but_a_real_drift_is_still_caught(self) -> None:
        import numpy as np

        from fastmdxplora.report.convergence import assess_run

        climbing = np.linspace(-490000, -400000, 50)
        out = assess_run({"potential_energy": climbing},
                         duration_ns=0.004, n_atoms=100)
        assert any("integration" in f for f in out["findings"])

    def test_the_report_says_too_short_rather_than_yes(self) -> None:
        import inspect

        from fastmdxplora.report import document

        source = inspect.getsource(document._convergence_section)
        assert "too short to say" in source


class TestFindingsAreSummarisedNotDumped:
    """The report printed two pages of raw tuples of atom indices.

    What an analysis worked out was recorded in ``options``, which is where
    settings live, and the report lists every setting. A binding-mode list and
    a transition matrix are structures somebody might want; they are not
    something anybody reads in a document.
    """

    def test_settings_and_findings_are_kept_apart(self, tmp_path) -> None:
        import mdtraj as md
        import numpy as np
        import pandas as pd

        from fastmdxplora.analysis.base import Analysis

        class _Learns(Analysis):
            name = "learns"
            description = "records what it found"

            def compute(self, traj):
                self.findings["what_it_worked_out"] = list(range(500))
                return pd.DataFrame({"frame": [0], "value": [1.0]})

            def plot(self, result, ax):
                ax.plot([0], [1])

        top = md.Topology()
        residue = top.add_residue("ALA", top.add_chain(), resSeq=1)
        top.add_atom("CA", md.element.carbon, residue)
        traj = md.Trajectory(
            xyz=np.zeros((1, 1, 3), dtype=np.float32), topology=top)

        analysis = _Learns(output_dir=str(tmp_path))
        analysis.run(traj)

        import json

        written = json.loads(
            next(tmp_path.rglob("options.json")).read_text(encoding="utf-8"))
        assert "what_it_worked_out" in written["findings"]
        assert "what_it_worked_out" not in written["options"]

    def test_the_interaction_analysis_records_them_as_findings(self) -> None:
        import inspect

        from fastmdxplora.analysis import pl_interactions

        source = inspect.getsource(pl_interactions)
        for name in ("ligand_chemistry", "binding_modes", "mode_transitions"):
            assert f'self.findings["{name}"]' in source, name
            assert f'self.options["{name}"]' not in source, name

    def test_the_report_says_what_they_amount_to(self) -> None:
        from fastmdxplora.report.document import _findings_notes

        notes = _findings_notes({
            "ligand_chemistry": {"source": "run", "charge_was_ambiguous": False},
            "binding_modes": [{"fraction": 0.42}],
            "mode_transitions": {"observed_transitions": 96, "supported": True},
        })
        joined = " ".join(notes)
        assert "resolved during setup" in joined
        assert "42%" in joined
        assert "96 changes" in joined
        # And nothing that would run to pages.
        assert all(len(note) < 200 for note in notes)

    def test_a_guessed_chemistry_is_named_as_a_guess(self) -> None:
        from fastmdxplora.report.document import _findings_notes

        notes = _findings_notes({"ligand_chemistry": {"source": "perceived"}})
        assert "guess" in " ".join(notes)

    def test_an_unsupportable_rate_is_reported_as_such(self) -> None:
        from fastmdxplora.report.document import _findings_notes

        notes = _findings_notes({
            "mode_transitions": {"observed_transitions": 2, "supported": False}})
        assert "too few for a rate" in " ".join(notes)

    def test_a_refused_measurement_is_mentioned(self) -> None:
        from fastmdxplora.report.document import _findings_notes

        notes = _findings_notes({"not_measured": {"salt_bridge": "charge unknown"}})
        assert "Not measured" in " ".join(notes)


class TestTheMethodsSectionReportsTheRestraints:
    """A restrained equilibration is part of a protocol. A reader cannot
    repeat a run that held the protein at a thousand kilojoules and released
    it in four steps unless the paragraph says so.
    """

    @staticmethod
    def _text(**extra):
        from pathlib import Path

        from fastmdxplora.report.methods import methods_paragraphs

        sim = {"parameters": dict({
            "timestep_fs": 2.0, "integrator": "langevin_middle",
            "temperature_K": 300.0, "friction_per_ps": 1.0,
            "nvt_steps": 25000, "npt_steps": 25000,
            "production_steps": 500000, "minimize": True,
        }, **extra), "platform_used": "CUDA", "pressure_bar_used": 1.0}
        return methods_paragraphs(Path("."), {}, sim)

    def test_what_was_held_and_how_hard(self) -> None:
        text = self._text(restrain="protein and not element H",
                          restraint_release=[1000, 500, 100, 0])
        assert "protein and not element H" in text
        assert "1000, 500, 100, 0" in text
        assert "kJ/mol/nm" in text

    def test_that_they_came_off_for_production(self) -> None:
        text = self._text(restrain="protein", restrain_production=False)
        assert "released before production" in text

    def test_and_that_they_did_not_where_they_did_not(self) -> None:
        """A biased trajectory reported as a free one is a claim the
        simulation does not support."""
        text = self._text(restrain="protein", restrain_production=True)
        assert "retained during production" in text
        assert "biased" in text

    def test_an_unrestrained_run_says_nothing_about_restraints(self) -> None:
        text = self._text()
        assert "restrained" not in text

    def test_a_clause_with_no_value_is_left_out(self) -> None:
        """It read "a friction coefficient of None ps⁻¹" when the number was
        absent, which is worse than omitting it: a reader cannot tell whether
        the run had no friction or the software lost it."""
        from pathlib import Path

        from fastmdxplora.report.methods import methods_paragraphs

        sim = {"parameters": {"timestep_fs": 2.0, "integrator": "verlet",
                              "temperature_K": 300.0, "production_steps": 1000}}
        text = methods_paragraphs(Path("."), {}, sim)
        assert "None" not in text


class TestTheReportDoesNotStateWhatDidNotHappen:
    """A methods section is the part of a report that gets published.

    Run on a tri-alanine in water, with no ligand in it, the methods paragraph
    read: "The ligand LIG was parameterized with openff-2.2.1." Both values
    behind that sentence are defaults --- LIG is the naming convention for a
    ligand if there is one, and the small-molecule force field is a property
    of the protein force field, present whether or not it was used --- so the
    sentence appeared in every run.
    """

    @staticmethod
    def _setup(**over):
        base = {
            "input": {"system": "/home/someone/ala3.pdb", "form": "pdb_file"},
            "parameters": {
                "ph": 7.4, "ligand": None, "ligand_name": "LIG",
                "ligand_forcefield": None, "solvent_padding_nm": 1.2,
                "box_shape": "cube", "water_model": "tip3p",
            },
            "resolved_forcefield": {
                "xmls": ["amber14-all.xml"], "water_model": "tip3p",
                "supports_ligand": True,
                "small_molecule_forcefield": "openff-2.2.1",
            },
            "n_atoms_solvated": 1271,
        }
        base["parameters"].update(over)
        return base

    def test_no_ligand_is_claimed_when_none_was_parameterised(
        self, tmp_path
    ) -> None:
        from fastmdxplora.report.methods import methods_paragraphs

        text = methods_paragraphs(tmp_path, self._setup(), {})
        assert "ligand" not in text.lower()

    def test_one_is_claimed_when_a_ligand_was(self, tmp_path) -> None:
        from fastmdxplora.report.methods import methods_paragraphs

        text = methods_paragraphs(
            tmp_path,
            self._setup(ligand=["ligands/BNZ.sdf"], ligand_name="BNZ",
                        ligand_forcefield="openff-2.2.1"),
            {})
        assert "BNZ" in text
        assert "openff-2.2.1" in text

    def test_the_coordinate_source_is_read_not_guessed(self, tmp_path) -> None:
        """A four-character filename is not a deposited structure. The origin
        was decided by the length of the string, so any such name was
        described as coming from the Protein Data Bank."""
        from fastmdxplora.report.methods import methods_paragraphs

        setup = self._setup()
        setup["input"] = {"system": "abcd", "form": "pdb_file"}
        assert "Protein Data Bank" not in methods_paragraphs(tmp_path, setup, {})

        setup["input"] = {"system": "181L", "form": "pdb_id"}
        assert "Protein Data Bank entry" in methods_paragraphs(tmp_path, setup, {})


class TestTheTitleNamesTheSystem:
    """The title read "FastMDXplora Study --- /home/claude/ala3.pdb": the
    author's home directory, on the first line of a document meant to be sent
    to somebody."""

    def test_a_path_becomes_a_name(self) -> None:
        from fastmdxplora.report.run import _system_label

        assert _system_label("/home/someone/ala3.pdb") == "ala3"
        assert _system_label("/Users/someone/runs/1UBQ.cif") == "1UBQ"

    def test_a_pdb_entry_stays_one(self) -> None:
        from fastmdxplora.report.run import _system_label

        assert _system_label("181L") == "181L"
        assert _system_label("1ubq") == "1UBQ"

    def test_two_suffixes_both_come_off(self) -> None:
        """Taking one left "trp_cage.pdb" as the name of a study."""
        from fastmdxplora.report.run import _system_label

        assert _system_label("trp_cage.pdb.gz") == "trp_cage"


class TestARequestedOutputThatCouldNotBeMadeIsRecorded:
    """The terminal warned that no PDF could be written and the manifest did
    not. Read from its files afterwards, the run showed four formats where
    five were asked for, with nothing to say the fifth had been attempted."""

    def test_the_reason_is_written_beside_the_outputs(self) -> None:
        import inspect

        # `fastmdxplora.report.run` is the function, re-exported over the
        # module of the same name.
        import fastmdxplora.report.run as report_run_module

        source = inspect.getsource(report_run_module)
        assert 'not_produced.append(("report.pdf", reason))' in source
        assert '"not_produced.json"' in source


class TestTheSummarySaysWhatTheStudyWas:
    """The first section a reader reads said "This report was generated
    automatically by FastMDXplora from the outputs of an end-to-end molecular
    dynamics study" --- a statement about the software, in a document about
    their system."""

    def test_it_states_the_system_and_the_run(self, tmp_path) -> None:
        import json

        from fastmdxplora.report.document import _study_in_one_paragraph

        (tmp_path / "setup").mkdir()
        (tmp_path / "setup" / "setup_parameters.json").write_text(
            json.dumps({"n_atoms_solvated": 37098}), encoding="utf-8")
        (tmp_path / "simulation").mkdir()
        (tmp_path / "simulation" / "simulation_parameters.json").write_text(
            json.dumps({"duration_ns_actual": 20.0,
                        "parameters": {"temperature_K": 300.0}}),
            encoding="utf-8")

        said = _study_in_one_paragraph(tmp_path)
        assert "37,098 atoms" in said
        assert "20 ns" in said and "300 K" in said

    def test_a_short_run_is_given_in_picoseconds(self, tmp_path) -> None:
        """"0.008 ns" is a number nobody says out loud."""
        import json

        from fastmdxplora.report.document import _study_in_one_paragraph

        (tmp_path / "simulation").mkdir()
        (tmp_path / "simulation" / "simulation_parameters.json").write_text(
            json.dumps({"duration_ns_actual": 0.008}), encoding="utf-8")
        assert "8 ps" in _study_in_one_paragraph(tmp_path)

    def test_nothing_recorded_means_nothing_claimed(self, tmp_path) -> None:
        """A sentence assembled from three absent values is worse than the
        generic one it replaces."""
        from fastmdxplora.report.document import _study_in_one_paragraph

        assert _study_in_one_paragraph(tmp_path) is None

    def test_the_convergence_counts_add_up(self) -> None:
        """A first version counted settled and unjudged from overlapping
        conditions, and reported three and six out of six."""
        import re

        from fastmdxplora.report.document import _what_the_run_supports

        text = _what_the_run_supports.__doc__
        assert text  # the behaviour is checked below against a real run

        import inspect
        source = inspect.getsource(_what_the_run_supports)
        assert "unjudged = total - settled - drifting" in source, (
            "the three states must be exclusive or the counts do not sum")


class TestTheSlidesArePresentable:
    """Twenty-one slides, twelve figures, not one number; three of them
    displaying filesystem paths from the machine that produced them; and the
    whole deck in 4:3."""

    @staticmethod
    def _recorded(tmp_path):
        import json

        (tmp_path / "setup").mkdir(exist_ok=True)
        (tmp_path / "setup" / "setup_parameters.json").write_text(json.dumps({
            "n_atoms_solvated": 37098,
            "parameters": {"ph": 7.4, "solvent_padding_nm": 1.0,
                           "box_shape": "cube", "ion_concentration_M": 0.15,
                           "ligand": ["ligands/BNZ.sdf"], "ligand_name": "BNZ",
                           "ligand_forcefield": "openff-2.2.1"},
            "resolved_forcefield": {"xmls": ["amber14-all.xml"],
                                    "water_model": "tip3p"},
        }), encoding="utf-8")
        (tmp_path / "simulation").mkdir(exist_ok=True)
        (tmp_path / "simulation" / "simulation_parameters.json").write_text(
            json.dumps({
                "duration_ns_actual": 20.0, "pressure_bar_used": 1.0,
                "platform_used": "CUDA",
                "parameters": {"temperature_K": 300.0, "timestep_fs": 2.0,
                               "integrator": "langevin_middle",
                               "constraints": "HBonds", "random_seed": 7},
            }), encoding="utf-8")

    def test_the_deck_is_widescreen(self) -> None:
        import inspect

        from fastmdxplora.report import slides

        source = inspect.getsource(slides._build_pptx)
        assert "Inches(13.333)" in source, "4:3 shows black bands on a projector"

    def test_the_setup_slide_states_the_system_not_a_path(self, tmp_path) -> None:
        from fastmdxplora.report.slides import _setup_bullets

        self._recorded(tmp_path)
        bullets = _setup_bullets(tmp_path)

        assert any("37,098 atoms" in b for b in bullets)
        assert any("amber14-all.xml" in b for b in bullets)
        assert any("BNZ" in b for b in bullets)
        assert not any("/" in b and b.startswith("  •") for b in bullets)

    def test_the_ensemble_comes_from_what_ran(self, tmp_path) -> None:
        """`pressure_bar` is what was asked for and is None when the default
        was taken; reading it called an NPT run at 1 bar an NVT one."""
        from fastmdxplora.report.slides import _simulation_bullets

        self._recorded(tmp_path)
        assert any("1 bar (NPT)" in b for b in _simulation_bullets(tmp_path))

    def test_the_seed_is_stated_either_way(self, tmp_path) -> None:
        from fastmdxplora.report.slides import _simulation_bullets

        self._recorded(tmp_path)
        assert any("Random seed: 7" in b for b in _simulation_bullets(tmp_path))

    def test_the_outline_does_not_also_point_at_the_json(self, tmp_path) -> None:
        """`list.extend` returns None, so "extend(...) or append(fallback)"
        ran the fallback every time: the outline carried the bullets and then
        told the reader to go and look at the manifest anyway."""
        import inspect

        from fastmdxplora.report import slides

        source = inspect.getsource(slides)
        assert ") or lines.append(" not in source


class TestTheBundleCarriesTheRunRecord:
    """The archive is the thing a study is sent as. It carried the outputs,
    including the trajectory, and not the manifest recording what produced
    them nor the configuration file that reproduces them.

    Not by exclusion: those two are written once every phase has finished, and
    the bundle is built during the report phase, so they did not exist yet.
    The module's own docstring claimed it contained "the top-level manifest".
    """

    def test_the_manifest_and_the_config_are_inside(self, tmp_path) -> None:
        import zipfile

        from fastmdxplora.orchestrator import FastMDXplora

        root = tmp_path / "run"
        (root / "report").mkdir(parents=True)
        with zipfile.ZipFile(root / "report" / "project_bundle.zip", "w") as z:
            z.writestr("analysis/rmsd/rmsd.dat", "0.1\n")
        (root / "manifest.json").write_text('{"phases": []}', encoding="utf-8")
        (root / "resolved_config.yml").write_text("output: run\n", encoding="utf-8")

        fmdx = FastMDXplora(system="1L2Y", output_dir=root)
        fmdx._add_run_record_to_bundle()

        inside = zipfile.ZipFile(root / "report" / "project_bundle.zip").namelist()
        assert "manifest.json" in inside
        assert "resolved_config.yml" in inside
        assert "analysis/rmsd/rmsd.dat" in inside, "nothing already there is lost"

    def test_it_is_added_after_they_are_written(self) -> None:
        """Which is the whole point: called earlier, there is nothing to add."""
        import inspect

        from fastmdxplora.orchestrator import FastMDXplora

        source = inspect.getsource(FastMDXplora)
        order = [source.index(call) for call in (
            "self._write_manifest()",
            "self._write_resolved_config()",
            "self._add_run_record_to_bundle()",
        )]
        assert order == sorted(order)

    def test_a_missing_bundle_is_not_an_error(self, tmp_path) -> None:
        """Bundles are optional, and a study that finished should not fail for
        want of an archive nobody asked for."""
        from fastmdxplora.orchestrator import FastMDXplora

        root = tmp_path / "run"
        root.mkdir()
        FastMDXplora(system="1L2Y", output_dir=root)._add_run_record_to_bundle()

    def test_nothing_is_added_twice(self, tmp_path) -> None:
        import zipfile

        from fastmdxplora.orchestrator import FastMDXplora

        root = tmp_path / "run"
        (root / "report").mkdir(parents=True)
        with zipfile.ZipFile(root / "report" / "project_bundle.zip", "w") as z:
            z.writestr("manifest.json", "{}")
        (root / "manifest.json").write_text('{"phases": []}', encoding="utf-8")

        fmdx = FastMDXplora(system="1L2Y", output_dir=root)
        fmdx._add_run_record_to_bundle()

        names = zipfile.ZipFile(root / "report" / "project_bundle.zip").namelist()
        assert names.count("manifest.json") == 1


class TestTheDashboardHeaderNamesTheSystem:
    """The same field the report and the slides show. The output folder and
    the `fastmdx gui --output ...` instruction keep their paths: that page is
    opened on the machine that produced the run, and both need one."""

    def test_the_header_shows_a_name(self) -> None:
        from fastmdxplora.gui.report_dashboard import _system_label

        assert _system_label("/home/someone/ala3.pdb") == "ala3"
        assert _system_label("181L") == "181L"
